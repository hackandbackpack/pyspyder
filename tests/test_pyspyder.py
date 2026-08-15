"""Regression tests for PySpyder.

Run with either:
    python -m unittest discover -s tests
    pytest tests

Stdlib unittest on purpose, so the suite needs nothing beyond what the tool
itself already requires. Office fixtures are generated into a temp directory at
run time rather than committed, so no binaries live in the tree.

Most of these exist because the behaviour they pin down was once wrong. Where
that is the case the test says so, so nobody "simplifies" it back.
"""

import io
import logging
import os
import shutil
import sys
import tempfile
import unittest
import zipfile
from datetime import datetime
from unittest import mock

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pyspyder as ps  # noqa: E402

logging.getLogger("pyspyder").setLevel(logging.CRITICAL)
# pypdf narrates every malformed document it is handed; the suite feeds it
# several on purpose. setup_logging() does the same for real runs.
logging.getLogger("pypdf").setLevel(logging.CRITICAL)


# ---------------------------------------------------------------------------
# Test doubles
# ---------------------------------------------------------------------------

class FakeResponse:
    def __init__(self, text="", status=200, ctype="text/html", url="",
                 body=b"", headers=None, fail_after=None):
        self.text = text
        self.status_code = status
        self.url = url
        self.body = body
        self.headers = {"Content-Type": ctype}
        if headers:
            self.headers.update(headers)
        self.fail_after = fail_after

    def raise_for_status(self):
        if self.status_code >= 400:
            raise RuntimeError(f"HTTP {self.status_code}")

    def close(self):
        pass

    def iter_content(self, chunk_size=8192):
        sent = 0
        for i in range(0, len(self.body), chunk_size):
            if self.fail_after is not None and sent >= self.fail_after:
                raise IOError("connection reset mid-transfer")
            chunk = self.body[i:i + chunk_size]
            sent += len(chunk)
            yield chunk


class FakeSession:
    """Routes URLs to canned responses and records every request made."""

    def __init__(self, routes=None, default=None):
        self.routes = routes or {}
        self.default = default if default is not None else FakeResponse("", 404)
        self.requested = []

    def head(self, url, **kwargs):
        return FakeResponse(url=url)

    def get(self, url, **kwargs):
        self.requested.append(url)
        route = self.routes.get(url, self.default)
        return route(url) if callable(route) else route

    def pages_requested(self):
        return [u for u in self.requested
                if "robots.txt" not in u and "sitemap" not in u]


def html(*links):
    body = "".join(f'<a href="{href}">x</a>' for href in links)
    return FakeResponse(f"<html><body>{body}</body></html>")


def document(body, ctype="application/pdf", **kwargs):
    """A response carrying file bytes rather than a page."""
    return FakeResponse(body=body, ctype=ctype, **kwargs)


def crawl(domain, session, file_types=("pdf",), **kwargs):
    """spider_domain with the delay stubbed out, returning (files, sleeps)."""
    kwargs.setdefault("max_depth", 2)
    kwargs.setdefault("delay", 0)
    kwargs.setdefault("ignore_robots", True)
    with mock.patch.object(ps.time, "sleep") as sleeper:
        found = ps.spider_domain(domain, list(file_types), session, **kwargs)
    return found, sleeper.call_count


# ---------------------------------------------------------------------------
# URL canonicalization
# ---------------------------------------------------------------------------

class TestUrlKey(unittest.TestCase):
    def test_collapses_equivalent_forms(self):
        canonical = ps.url_key("https://example.com/a.pdf")
        for variant in [
            "https://EXAMPLE.com/a.pdf",
            "http://example.com/a.pdf",          # scheme is not part of identity
            "https://example.com/a.pdf#frag",
            "https://example.com/a.pdf?utm_source=x&fbclid=y",
        ]:
            self.assertEqual(ps.url_key(variant), canonical, variant)

    def test_distinguishes_different_resources(self):
        keys = {ps.url_key(u) for u in [
            "https://example.com/a.pdf",
            "https://example.com/b.pdf",
            "https://example.com/A.pdf",         # path case is significant
            "https://other.com/a.pdf",
            "https://example.com/a.pdf?id=1",
        ]}
        self.assertEqual(len(keys), 5)

    def test_query_order_does_not_matter(self):
        self.assertEqual(ps.url_key("https://example.com/x?b=2&a=1"),
                         ps.url_key("https://example.com/x?a=1&b=2"))

    def test_meaningful_query_is_kept(self):
        self.assertIn("file=secret.pdf",
                      ps.url_key("https://example.com/dl.php?file=secret.pdf"))

    def test_malformed_port_does_not_raise(self):
        # parsed.port raises ValueError on these, which once aborted a whole
        # crawl through main()'s exception handler.
        for bad in ["https://example.com:99999/a.pdf",
                    "https://example.com:abc/a.pdf",
                    "https://example.com:-1/a.pdf"]:
            self.assertIn("example.com", ps.url_key(bad), bad)

    def test_userinfo_is_not_part_of_identity(self):
        self.assertEqual(ps.url_key("https://user:pw@example.com/a.pdf"),
                         ps.url_key("https://example.com/a.pdf"))


class TestDedupeUrls(unittest.TestCase):
    def test_returns_urls_verbatim(self):
        # The canonical key rewrites the query; what we hand back must not be
        # rewritten, because it is what gets requested.
        for url in ["https://example.com/f.pdf?flag",
                    "https://example.com/f.pdf?token=a/b+c",
                    "https://example.com/f.pdf?a=%20space"]:
            self.assertEqual(ps.dedupe_urls([url]), [url])

    def test_prefers_https_over_http(self):
        self.assertEqual(
            ps.dedupe_urls(["http://example.com/a.pdf", "https://example.com/a.pdf"]),
            ["https://example.com/a.pdf"])

    def test_keeps_http_only_resources(self):
        result = ps.dedupe_urls(["http://example.com/only.pdf"])
        self.assertEqual(result, ["http://example.com/only.pdf"])

    def test_is_deterministic_regardless_of_input_order(self):
        urls = ["https://example.com/b.pdf", "https://example.com/a.pdf",
                "https://example.com/a.pdf?utm_source=1"]
        self.assertEqual(ps.dedupe_urls(urls), ps.dedupe_urls(list(reversed(urls))))


# ---------------------------------------------------------------------------
# Scope
# ---------------------------------------------------------------------------

class TestScope(unittest.TestCase):
    def test_target_and_subdomains_are_in_scope(self):
        for host in ["example.com", "www.example.com", "cdn.a.example.com",
                     "EXAMPLE.COM", "example.com:8443", "example.com."]:
            self.assertTrue(ps.host_in_scope(host, {"example.com"}), host)

    def test_lookalike_hosts_are_out_of_scope(self):
        # A substring test accepted all of these, which is how a spider ends up
        # crawling somebody else's estate.
        for host in ["example.com.attacker.test", "notexample.com",
                     "example.community", "attacker.test", ""]:
            self.assertFalse(ps.host_in_scope(host, {"example.com"}), host)

    def test_ipv6_hosts_are_compared_properly(self):
        # split(":")[0] reduced every IPv6 literal to "[", so any two of them
        # compared equal.
        self.assertFalse(ps.host_in_scope("[2001:db8::2]", {"[2001:db8::1]"}))
        self.assertTrue(ps.host_in_scope("[2001:db8::1]:9999", {"[2001:db8::1]"}))

    def test_bare_host_strips_port_and_userinfo(self):
        self.assertEqual(ps.bare_host("user@example.com:8443"), "example.com")
        self.assertEqual(ps.bare_host("[::1]:8443"), "::1")
        self.assertEqual(ps.bare_host(""), "")

    def test_normalize_domain_accepts_what_users_type(self):
        cases = {
            "example.com": ("example.com", ""),
            "https://example.com": ("example.com", ""),
            "EXAMPLE.com/docs/": ("example.com", "docs"),
            "  example.com  ": ("example.com", ""),
            "http://user@example.com/x/y": ("example.com", "x/y"),
            "example.com:8443/": ("example.com:8443", ""),
        }
        for raw, expected in cases.items():
            self.assertEqual(ps.normalize_domain(raw), expected, raw)


class TestSession(unittest.TestCase):
    def test_advertises_only_decodable_encodings(self):
        # The session once advertised "br" while no brotli decoder was installed.
        # requests handed back the still-compressed bytes, every page parsed to
        # zero links, and the crawl reported nothing found - a silent miss on any
        # Brotli server. Pin the invariant: never claim an encoding we can't
        # decode. Stays green if brotli is later added as a dependency, because
        # CONTENT_DECODERS grows to include it.
        from urllib3.response import HTTPResponse

        decodable = set(HTTPResponse.CONTENT_DECODERS) | {"identity", "*"}
        advertised = ps.get_session().headers["Accept-Encoding"]
        for coding in advertised.split(","):
            coding = coding.split(";")[0].strip().lower()
            if coding:
                self.assertIn(coding, decodable, advertised)


class TestSeedFailure(unittest.TestCase):
    def test_blocked_seed_url_warns_loudly(self):
        # A 403 on the starting page is an access block, not an empty site, and
        # was being swallowed into a quiet "Found 0 files" - the operator would
        # conclude the target exposes nothing.
        session = FakeSession({
            "https://example.com/": FakeResponse(status=403, ctype="text/html"),
        })
        with self.assertLogs("pyspyder", level="WARNING") as cm:
            found, _ = crawl("example.com", session)
        self.assertEqual(found, set())
        message = " ".join(cm.output).lower()
        self.assertIn("403", message)
        self.assertIn("--user-agent", message)

    def test_reachable_seed_does_not_warn(self):
        # The block warning must not fire on an ordinary, working crawl.
        session = FakeSession({
            "https://example.com/": html("https://example.com/a.pdf"),
        })
        logger = logging.getLogger("pyspyder")
        captured = []
        handler = logging.Handler()
        handler.emit = captured.append
        old_level = logger.level
        logger.setLevel(logging.DEBUG)
        logger.addHandler(handler)
        try:
            found, _ = crawl("example.com", session)
        finally:
            logger.removeHandler(handler)
            logger.setLevel(old_level)
        self.assertIn("https://example.com/a.pdf", found)
        warnings = [r.getMessage() for r in captured if r.levelno >= logging.WARNING]
        self.assertEqual(warnings, [])


class TestSpiderScope(unittest.TestCase):
    def test_lookalike_host_is_neither_crawled_nor_collected(self):
        session = FakeSession({
            "https://example.com/": html(
                "https://example.com/real.pdf",
                "https://cdn.example.com/legit.pdf",
                "https://example.com.attacker.test/evil.pdf",
                "https://example.com.attacker.test/page.html",
                "https://notexample.com/evil2.pdf"),
            "https://example.com.attacker.test/page.html": html("/pivoted.pdf"),
        })
        found, _ = crawl("example.com", session)
        self.assertEqual(found, {"https://example.com/real.pdf",
                                 "https://cdn.example.com/legit.pdf"})
        self.assertNotIn("https://example.com.attacker.test/page.html",
                         session.requested)

    def test_non_http_schemes_are_ignored(self):
        session = FakeSession({
            "https://example.com/": html("mailto:a@example.com", "javascript:void(0)",
                                         "ftp://example.com/a.pdf", "/real.pdf"),
        })
        found, _ = crawl("example.com", session)
        self.assertEqual(found, {"https://example.com/real.pdf"})

    def test_offsite_sitemap_from_robots_is_not_fetched(self):
        session = FakeSession({
            "https://example.com/robots.txt": FakeResponse(
                "User-agent: *\nSitemap: https://attacker.test/sitemap.xml\n",
                ctype="text/plain"),
            "https://example.com/": html(),
        })
        crawl("example.com", session, ignore_robots=False)
        self.assertNotIn("https://attacker.test/sitemap.xml", session.requested)

    def test_offsite_sitemap_entries_are_dropped(self):
        session = FakeSession({
            "https://example.com/": html(),
            "https://example.com/sitemap.xml": FakeResponse(
                '<?xml version="1.0"?><urlset>'
                '<url><loc>https://example.com/ours.pdf</loc></url>'
                '<url><loc>https://elsewhere.test/theirs.pdf</loc></url>'
                '</urlset>', ctype="text/xml"),
        })
        found, _ = crawl("example.com", session)
        self.assertEqual(found, {"https://example.com/ours.pdf"})

    def test_page_redirected_offsite_is_not_parsed(self):
        session = FakeSession({
            "https://example.com/": html("/redirector.html"),
            "https://example.com/redirector.html": FakeResponse(
                '<a href="/their-secret.pdf">x</a>',
                url="https://thirdparty.test/landing.html"),
        })
        found, _ = crawl("example.com", session)
        self.assertEqual(found, set())


# ---------------------------------------------------------------------------
# robots.txt
# ---------------------------------------------------------------------------

class TestRobots(unittest.TestCase):
    def _session(self):
        return FakeSession({
            "https://example.com/robots.txt": FakeResponse(
                "User-agent: *\nDisallow: /private/\n", ctype="text/plain"),
            "https://example.com/": html("/private/blocked.pdf", "/public/ok.pdf",
                                         "/private/page.html"),
            "https://example.com/private/page.html": html("/private/deep.pdf"),
        })

    def test_disallowed_files_are_not_collected(self):
        # The gate used to cover crawled pages only, so disallowed documents
        # were downloaded anyway.
        found, _ = crawl("example.com", self._session(), ignore_robots=False)
        self.assertEqual(found, {"https://example.com/public/ok.pdf"})

    def test_disallowed_pages_are_not_crawled(self):
        session = self._session()
        crawl("example.com", session, ignore_robots=False)
        self.assertNotIn("https://example.com/private/page.html", session.requested)

    def test_ignore_robots_collects_everything(self):
        found, _ = crawl("example.com", self._session(), ignore_robots=True)
        self.assertIn("https://example.com/private/blocked.pdf", found)

    def test_rules_are_per_host(self):
        # One parser for the base host used to be applied to every subdomain, so
        # www's Disallow silently suppressed cdn's files.
        session = FakeSession({
            "https://example.com/robots.txt": FakeResponse(
                "User-agent: *\nDisallow: /docs/\n", ctype="text/plain"),
            "https://cdn.example.com/robots.txt": FakeResponse(
                "User-agent: *\nAllow: /\n", ctype="text/plain"),
            "https://example.com/": html("https://cdn.example.com/docs/theirs.pdf",
                                         "/docs/ours.pdf"),
        })
        found, _ = crawl("example.com", session, ignore_robots=False)
        self.assertEqual(found, {"https://cdn.example.com/docs/theirs.pdf"})

    def test_sitemap_directive_is_harvested(self):
        session = FakeSession({
            "https://example.com/robots.txt": FakeResponse(
                "User-agent: *\nSitemap: https://example.com/sm/extra.xml\n",
                ctype="text/plain"),
            "https://example.com/": html(),
            "https://example.com/sm/extra.xml": FakeResponse(
                '<?xml version="1.0"?><urlset><url>'
                '<loc>https://example.com/from-robots.pdf</loc></url></urlset>',
                ctype="text/xml"),
        })
        found, _ = crawl("example.com", session, ignore_robots=False)
        self.assertIn("https://example.com/from-robots.pdf", found)

    def test_missing_robots_allows_everything(self):
        session = FakeSession({"https://example.com/": html("/a.pdf")})
        found, _ = crawl("example.com", session, ignore_robots=False)
        self.assertEqual(found, {"https://example.com/a.pdf"})

    def test_unreachable_robots_allows_everything(self):
        session = FakeSession({
            "https://example.com/robots.txt": FakeResponse("", 500, ctype="text/plain"),
            "https://example.com/": html("/a.pdf"),
        })
        found, _ = crawl("example.com", session, ignore_robots=False)
        self.assertEqual(found, {"https://example.com/a.pdf"})

    def test_disallow_all_blocks_and_is_reported(self):
        session = FakeSession({
            "https://example.com/robots.txt": FakeResponse(
                "User-agent: *\nDisallow: /\n", ctype="text/plain"),
            "https://example.com/": html("/a.pdf"),
        })
        policy = ps.RobotsPolicy(session)
        self.assertFalse(policy.allowed("https://example.com/a.pdf"))
        self.assertEqual(policy.blocked, 1)

    def test_robots_is_fetched_once_per_host(self):
        session = FakeSession({
            "https://example.com/robots.txt": FakeResponse("", 404, ctype="text/plain")})
        policy = ps.RobotsPolicy(session)
        for _ in range(5):
            policy.allowed("https://example.com/a.pdf")
        self.assertEqual(session.requested.count("https://example.com/robots.txt"), 1)

    def test_relative_sitemap_directive_is_resolved(self):
        session = FakeSession({
            "https://example.com/robots.txt": FakeResponse(
                "User-agent: *\nSitemap: /sm/relative.xml\n", ctype="text/plain"),
            "https://example.com/": html(),
            "https://example.com/sm/relative.xml": FakeResponse(
                '<?xml version="1.0"?><urlset><url>'
                '<loc>https://example.com/relative-found.pdf</loc></url></urlset>',
                ctype="text/xml"),
        })
        found, _ = crawl("example.com", session, ignore_robots=False)
        self.assertIn("https://example.com/relative-found.pdf", found)


# ---------------------------------------------------------------------------
# Crawling
# ---------------------------------------------------------------------------

class TestCrawl(unittest.TestCase):
    def test_depth_limit(self):
        routes = {
            "https://example.com/": html("/one.html"),
            "https://example.com/one.html": html("/two.html", "/depth1.pdf"),
            "https://example.com/two.html": html("/depth2.pdf"),
        }
        found, _ = crawl("example.com", FakeSession(routes), max_depth=1)
        self.assertEqual(found, {"https://example.com/depth1.pdf"})

    def test_max_pages_is_honored(self):
        links = [f"/p{i}.html" for i in range(30)]
        routes = {"https://example.com/": html(*links)}
        for link in links:
            routes[f"https://example.com{link}"] = html()
        session = FakeSession(routes)
        crawl("example.com", session, max_pages=5)
        self.assertLessEqual(len(session.pages_requested()), 5)

    def test_tracking_variants_fetched_once(self):
        routes = {
            "https://example.com/": html("/page.html", "/page.html?utm_source=a",
                                         "/page.html?fbclid=b", "/page.html#frag"),
            "https://example.com/page.html": html(),
        }
        session = FakeSession(routes)
        crawl("example.com", session)
        self.assertEqual(len(session.pages_requested()), 2)

    def test_relative_links_resolve_against_the_redirect_target(self):
        # /docs 301s to /docs/. Joining onto the requested url gave
        # /budget.xlsx, which does not exist.
        session = FakeSession({
            "https://example.com/": html("/docs"),
            "https://example.com/docs": FakeResponse(
                '<a href="budget.xlsx">x</a>', url="https://example.com/docs/"),
        })
        found, _ = crawl("example.com", session, file_types=("xlsx",))
        self.assertEqual(found, {"https://example.com/docs/budget.xlsx"})

    def test_base_href_overrides_the_document_url(self):
        session = FakeSession({
            "https://example.com/": FakeResponse(
                '<html><head><base href="/files/"></head>'
                '<body><a href="report.pdf">x</a></body></html>'),
        })
        found, _ = crawl("example.com", session, max_depth=1)
        self.assertEqual(found, {"https://example.com/files/report.pdf"})

    def test_xhtml_is_parsed(self):
        session = FakeSession({
            "https://example.com/": FakeResponse(
                '<a href="/a.pdf">x</a>', ctype="application/xhtml+xml"),
        })
        found, _ = crawl("example.com", session)
        self.assertEqual(found, {"https://example.com/a.pdf"})

    def test_non_html_is_not_parsed(self):
        session = FakeSession({
            "https://example.com/": FakeResponse(
                '<a href="/a.pdf">x</a>', ctype="application/octet-stream"),
        })
        found, _ = crawl("example.com", session)
        self.assertEqual(found, set())

    def test_delay_applies_after_failed_and_non_html_fetches(self):
        # The sleep used to sit at the bottom of the loop, so a 404 skipped it
        # even though the request had already gone out.
        links = [f"/p{i}.html" for i in range(4)]
        base = {"https://example.com/": html(*links)}

        outcomes = {
            "ok": FakeResponse("<html></html>"),
            "404": FakeResponse("", 404),
            "non-html": FakeResponse("data", ctype="application/octet-stream"),
        }
        counts = {}
        for label, leaf in outcomes.items():
            routes = dict(base)
            for link in links:
                routes[f"https://example.com{link}"] = leaf
            session = FakeSession(routes)
            _, sleeps = crawl("example.com", session)
            counts[label] = (len(session.pages_requested()), sleeps)

        self.assertEqual(len({c for c in counts.values()}), 1, counts)

    def test_requests_the_url_as_discovered(self):
        odd = "https://example.com/docs/f%20i+le.pdf?flag&token=a/b+c"
        session = FakeSession({"https://example.com/": html(odd)})
        found, _ = crawl("example.com", session)
        self.assertEqual(found, {odd})
        self.assertEqual(ps.dedupe_urls(found), [odd])


# ---------------------------------------------------------------------------
# Sitemaps
# ---------------------------------------------------------------------------

class TestSitemap(unittest.TestCase):
    def _parse(self, session, **kwargs):
        kwargs.setdefault("scope_hosts", {"example.com"})
        kwargs.setdefault("robots", ps.RobotsPolicy(session, ignore=True))
        kwargs.setdefault("delay", 0)
        with mock.patch.object(ps.time, "sleep") as sleeper:
            found = ps._parse_sitemap("https://example.com", ["pdf"], session, **kwargs)
        return found, sleeper.call_count

    def test_nested_sitemap_index_is_followed(self):
        session = FakeSession({
            "https://example.com/sitemap.xml": FakeResponse(
                '<?xml version="1.0"?><sitemapindex><sitemap>'
                '<loc>https://example.com/sm2.xml</loc></sitemap></sitemapindex>',
                ctype="text/xml"),
            "https://example.com/sm2.xml": FakeResponse(
                '<?xml version="1.0"?><urlset><url>'
                '<loc>https://example.com/nested.pdf</loc></url></urlset>',
                ctype="text/xml"),
        })
        found, _ = self._parse(session)
        self.assertEqual(found, {"https://example.com/nested.pdf"})

    def test_unusable_bodies_do_not_raise(self):
        for body in ["", "<!DOCTYPE html><html><title>404</title>", "\x00\x01 junk",
                     "<urlset><url><loc>https://example.com/a.pdf</lo"]:
            session = FakeSession(default=FakeResponse(body, ctype="text/xml"))
            found, _ = self._parse(session)
            self.assertIsInstance(found, set)

    def test_missing_xml_parser_does_not_sink_the_crawl(self):
        # A FeatureNotFound used to propagate out of spider_domain, and main()
        # discarded every file the crawl had already found.
        session = FakeSession({
            "https://example.com/": html("/found-by-crawl.pdf"),
            "https://example.com/sitemap.xml": FakeResponse("<urlset/>", ctype="text/xml"),
        })
        real = ps.BeautifulSoup

        def no_xml(markup, features=None, *a, **kw):
            if features == "xml":
                raise ps.bs4.FeatureNotFound("no xml tree builder")
            return real(markup, features, *a, **kw)

        with mock.patch.object(ps, "BeautifulSoup", no_xml):
            found, _ = crawl("example.com", session)
        self.assertEqual(found, {"https://example.com/found-by-crawl.pdf"})

    def test_sitemap_count_is_capped(self):
        children = [f"https://example.com/sm{i}.xml" for i in range(ps.MAX_SITEMAPS * 2)]
        index = "".join(f"<sitemap><loc>{c}</loc></sitemap>" for c in children)
        routes = {"https://example.com/sitemap.xml": FakeResponse(
            f'<?xml version="1.0"?><sitemapindex>{index}</sitemapindex>', ctype="text/xml")}
        for child in children:
            routes[child] = FakeResponse('<?xml version="1.0"?><urlset/>', ctype="text/xml")
        session = FakeSession(routes)
        self._parse(session)
        self.assertLessEqual(len(session.requested), ps.MAX_SITEMAPS)

    def test_delay_is_observed_between_sitemap_fetches(self):
        session = FakeSession({
            "https://example.com/sitemap.xml": FakeResponse(
                '<?xml version="1.0"?><urlset/>', ctype="text/xml")})
        _, sleeps = self._parse(session, delay=1)
        self.assertEqual(sleeps, 1)


# ---------------------------------------------------------------------------
# Downloading
# ---------------------------------------------------------------------------

PDF_BYTES = b"%PDF-1.4\n" + b"payload" * 20 + b"\n%%EOF\n"


class TestDownload(unittest.TestCase):
    def setUp(self):
        self.out = tempfile.mkdtemp(prefix="pyspyder-dl-")
        self.addCleanup(shutil.rmtree, self.out, ignore_errors=True)

    def _download(self, urls, session, **kwargs):
        kwargs.setdefault("delay", 0)
        with mock.patch.object(ps.time, "sleep") as sleeper:
            got = ps.download_files(urls, self.out, session, **kwargs)
        return got, sleeper.call_count

    def test_distinct_files_with_one_basename(self):
        # Three urls, two sharing a basename and a third genuinely named
        # report_1.pdf. This once reported 3 of 3 while writing 2 files and
        # attributed the third url to an unrelated file.
        urls = ["https://example.com/a/report.pdf",
                "https://example.com/b/report.pdf",
                "https://example.com/c/report_1.pdf"]
        bodies = {u: document(b"%PDF-1.4 " + tag)
                  for u, tag in zip(urls, [b"AAA", b"BBB", b"CCC"])}
        got, _ = self._download(urls, FakeSession(bodies))

        self.assertEqual(len(got), 3)
        self.assertEqual(len(os.listdir(self.out)), 3)
        for url, path in got:
            tag = {"a": b"AAA", "b": b"BBB", "c": b"CCC"}[url.split("/")[-2]]
            with open(path, "rb") as fh:
                self.assertIn(tag, fh.read(), f"{url} maps to the wrong file")

    def test_interrupted_transfer_leaves_nothing(self):
        url = "https://example.com/big.pdf"
        big = b"%PDF-1.4\n" + b"z" * 40000
        got, _ = self._download(
            [url], FakeSession({url: document(big, fail_after=16384)}))
        self.assertEqual(got, [])
        self.assertEqual(os.listdir(self.out), [])

    def test_retry_after_interruption_produces_a_whole_file(self):
        url = "https://example.com/big.pdf"
        big = b"%PDF-1.4\n" + b"z" * 40000
        self._download([url], FakeSession({url: document(big, fail_after=16384)}))
        got, _ = self._download([url], FakeSession({url: document(big)}))
        self.assertEqual(len(got), 1)
        self.assertEqual(os.path.getsize(got[0][1]), len(big))

    def test_document_url_serving_html_is_skipped(self):
        url = "https://example.com/login.pdf"
        got, _ = self._download([url], FakeSession(
            {url: document(b"<html>login</html>", ctype="text/html")}))
        self.assertEqual(got, [])
        self.assertEqual(os.listdir(self.out), [])

    def test_size_limit_from_content_length(self):
        url = "https://example.com/huge.pdf"
        response = document(PDF_BYTES, headers={"Content-Length": "999999999"})
        got, _ = self._download([url], FakeSession({url: response}), max_bytes=1000)
        self.assertEqual(got, [])

    def test_size_limit_while_streaming_when_header_lies(self):
        url = "https://example.com/huge.pdf"
        response = document(b"%PDF-1.4" + b"x" * 50000,
                            headers={"Content-Length": "10"})
        got, _ = self._download([url], FakeSession({url: response}), max_bytes=1000)
        self.assertEqual(got, [])
        self.assertEqual(os.listdir(self.out), [])

    def test_delay_applies_to_skipped_but_fetched_files(self):
        # A request went out in each of these cases, so the throttle is owed.
        urls = [f"https://example.com/a{i}.pdf" for i in range(3)]
        for response in [document(b"<html>", ctype="text/html"),
                         document(PDF_BYTES,
                                  headers={"Content-Length": "999999999"}),
                         FakeResponse("", 404)]:
            with self.subTest(response=response.status_code):
                shutil.rmtree(self.out, ignore_errors=True)
                os.makedirs(self.out)
                _, sleeps = self._download(
                    urls, FakeSession({u: response for u in urls}),
                    delay=1, max_bytes=1000)
                self.assertEqual(sleeps, len(urls))

    def test_no_delay_when_nothing_was_requested(self):
        url = "https://example.com/a.pdf"
        with open(os.path.join(self.out, "a.pdf"), "wb") as fh:
            fh.write(PDF_BYTES)
        got, sleeps = self._download([url], FakeSession(), delay=1)
        self.assertEqual(len(got), 1)
        self.assertEqual(sleeps, 0)

    def test_skipped_download_releases_its_filename(self):
        urls = ["https://example.com/a/report.pdf", "https://example.com/b/report.pdf"]
        session = FakeSession({
            urls[0]: document(b"<html>", ctype="text/html"),
            urls[1]: document(PDF_BYTES),
        })
        got, _ = self._download(urls, session)
        self.assertEqual([os.path.basename(p) for _, p in got], ["report.pdf"])

    def test_extension_content_mismatch_is_kept_but_flagged(self):
        url = "https://example.com/notreally.pdf"
        with self.assertLogs("pyspyder", level="WARNING") as logs:
            got, _ = self._download(
                [url], FakeSession({url: document(b"GIF89a not a pdf")}))
        self.assertEqual(len(got), 1)
        self.assertTrue(any("corrupt" in line for line in logs.output))


class TestFilenames(unittest.TestCase):
    def test_path_traversal_cannot_escape(self):
        for url in ["https://example.com/a/..%2f..%2f..%2fWindows%2fevil.pdf",
                    "https://example.com/%2e%2e%5c%2e%2e%5cevil.pdf"]:
            name = ps._extract_filename(url)
            self.assertNotIn("..", name)
            self.assertNotIn("/", name)
            self.assertNotIn("\\", name)

    def test_windows_device_names_are_escaped(self):
        # CON.pdf opens a console device instead of creating a file.
        self.assertEqual(ps._extract_filename("https://example.com/CON.pdf"), "_CON.pdf")
        self.assertEqual(ps._extract_filename("https://example.com/lpt1.pdf"), "_lpt1.pdf")

    def test_query_string_is_not_part_of_the_name(self):
        self.assertEqual(
            ps._extract_filename("https://example.com/r.pdf?v=2&t=1"), "r.pdf")

    def test_percent_encoding_is_decoded(self):
        self.assertEqual(
            ps._extract_filename("https://example.com/q%20report.pdf"), "q report.pdf")

    def test_unicode_survives(self):
        self.assertEqual(ps._extract_filename("https://example.com/résumé.pdf"),
                         "résumé.pdf")

    def test_long_names_are_capped_including_the_extension(self):
        for url in ["https://example.com/" + "b" * 300 + ".pdf",
                    "https://example.com/x." + "a" * 300]:
            name = ps._extract_filename(url)
            self.assertLessEqual(len(name),
                                 ps.MAX_FILENAME_LENGTH + ps.MAX_EXTENSION_LENGTH, url)

    def test_extension_is_preserved_for_dispatch(self):
        # ".pdf" splits to name=".pdf", ext="", which lost the extension the
        # metadata extractors dispatch on.
        name = ps._extract_filename("https://example.com/%20%20.pdf")
        self.assertTrue(name.endswith(".pdf"), name)

    def test_directory_urls_get_a_placeholder(self):
        self.assertEqual(ps._extract_filename("https://example.com/"), "unknown_file")


# ---------------------------------------------------------------------------
# Metadata
# ---------------------------------------------------------------------------

class TestMetadataExtraction(unittest.TestCase):
    """Real files, generated at run time, with values we can assert on."""

    @classmethod
    def setUpClass(cls):
        cls.dir = tempfile.mkdtemp(prefix="pyspyder-fixtures-")

        from pypdf import PdfWriter
        writer = PdfWriter()
        writer.add_blank_page(width=100, height=100)
        writer.add_metadata({
            "/Author": "jsmith",
            "/Creator": "Microsoft Word 2019",
            "/Producer": "Acrobat Distiller 21.0",
            "/Title": "Annual Report",
            "/Subject": "contact admin@example.com",
        })
        cls.pdf = os.path.join(cls.dir, "report.pdf")
        with open(cls.pdf, "wb") as fh:
            writer.write(fh)

        from docx import Document
        document = Document()
        document.add_paragraph("body")
        props = document.core_properties
        props.author = "jane.doe"
        props.last_modified_by = "admin"
        props.title = "Org Plan"
        props.created = datetime(2023, 11, 1, 9, 0, 0)
        cls.docx = os.path.join(cls.dir, "plan.docx")
        document.save(cls.docx)
        _inject_app_xml(cls.docx)

        from openpyxl import Workbook
        workbook = Workbook()
        workbook.active["A1"] = "data"
        workbook.properties.creator = "bob.jones"
        workbook.properties.lastModifiedBy = "svc_backup"
        cls.xlsx = os.path.join(cls.dir, "budget.xlsx")
        workbook.save(cls.xlsx)

        from pptx import Presentation
        deck = Presentation()
        deck.slides.add_slide(deck.slide_layouts[6])
        deck.core_properties.author = "carol.white"
        deck.core_properties.last_modified_by = "tmiller"
        cls.pptx = os.path.join(cls.dir, "deck.pptx")
        deck.save(cls.pptx)

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls.dir, ignore_errors=True)

    def test_pdf(self):
        meta = ps.extract_metadata(self.pdf)
        self.assertEqual(meta["Author"], "jsmith")
        self.assertEqual(meta["Creator"], "Microsoft Word 2019")
        self.assertEqual(meta["Title"], "Annual Report")

    def test_docx(self):
        meta = ps.extract_metadata(self.docx)
        self.assertEqual(meta["Author"], "jane.doe")
        self.assertEqual(meta["LastModifiedBy"], "admin")

    def test_docx_extended_properties(self):
        # Manager/Company/Template live in docProps/app.xml, which the core
        # properties APIs do not expose.
        meta = ps.extract_metadata(self.docx)
        self.assertEqual(meta["Manager"], "rmanager")
        self.assertEqual(meta["Company"], "Acme Corporation")
        self.assertEqual(meta["Template"], "CorpTemplate.dotm")

    def test_xlsx(self):
        meta = ps.extract_metadata(self.xlsx)
        self.assertEqual(meta["Author"], "bob.jones")
        self.assertEqual(meta["LastModifiedBy"], "svc_backup")

    def test_pptx(self):
        meta = ps.extract_metadata(self.pptx)
        self.assertEqual(meta["Author"], "carol.white")
        self.assertEqual(meta["LastModifiedBy"], "tmiller")

    def test_corrupt_file_costs_one_warning_not_the_run(self):
        broken = os.path.join(self.dir, "broken.pdf")
        with open(broken, "wb") as fh:
            fh.write(b"<html>not a pdf</html>")
        with self.assertLogs("pyspyder", level="WARNING"):
            self.assertEqual(ps.extract_metadata(broken), {})

    def test_unknown_extension_yields_nothing(self):
        other = os.path.join(self.dir, "notes.txt")
        with open(other, "w") as fh:
            fh.write("hello")
        self.assertEqual(ps.extract_metadata(other), {})

    def test_oversized_app_xml_is_not_read(self):
        # Deflate reaches ~1000:1, so a small entry can expand to hundreds of MB.
        bomb = os.path.join(self.dir, "bomb.docx")
        with zipfile.ZipFile(bomb, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("docProps/app.xml",
                        b"<Properties><Company>x</Company></Properties>"
                        + b" " * (ps.MAX_APP_XML_BYTES * 4))
        meta = {}
        ps._extract_ooxml_extended(bomb, meta)
        self.assertEqual(meta, {})

    def test_metadata_only_walks_subdirectories(self):
        nested = os.path.join(self.dir, "sub", "deeper")
        os.makedirs(nested, exist_ok=True)
        shutil.copy(self.pdf, os.path.join(nested, "copy.pdf"))
        names = [n for n, _ in ps.find_local_files(self.dir, ["pdf"])]
        self.assertTrue(any("copy.pdf" in n for n in names), names)

    def test_file_type_filter_accepts_a_leading_dot(self):
        self.assertTrue(ps.find_local_files(self.dir, [".pdf"]))


def _inject_app_xml(path):
    app_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/'
        '2006/extended-properties">'
        "<Application>Microsoft Office Word</Application>"
        "<AppVersion>16.0000</AppVersion>"
        "<Company>Acme Corporation</Company>"
        "<Manager>rmanager</Manager>"
        "<Template>CorpTemplate.dotm</Template>"
        "</Properties>"
    )
    buffer = io.BytesIO()
    with zipfile.ZipFile(path) as src:
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as dst:
            for item in src.infolist():
                if item.filename != "docProps/app.xml":
                    dst.writestr(item, src.read(item.filename))
            dst.writestr("docProps/app.xml", app_xml)
    with open(path, "wb") as fh:
        fh.write(buffer.getvalue())


class TestOleText(unittest.TestCase):
    def test_bytes_and_str_are_both_handled(self):
        # VT_LPSTR arrives as bytes, VT_LPWSTR as str. An unconditional
        # .decode() raised AttributeError, and because extract_metadata catches
        # everything that discarded the whole file's metadata.
        self.assertEqual(ps._ole_text(b"bytes.user"), "bytes.user")
        self.assertEqual(ps._ole_text("str.user"), "str.user")

    def test_undecodable_bytes_do_not_raise(self):
        self.assertIsInstance(ps._ole_text(b"\xff\xfe bad"), str)

    def test_extractor_keeps_str_fields(self):
        class Meta:
            author = "str.author"
            last_saved_by = b"bytes.user"
            company = title = subject = None
            creating_application = manager = keywords = None
            create_time = last_saved_time = None

        class Ole:
            def get_metadata(self):
                return Meta()

            def close(self):
                pass

        # olefile is imported inside the extractor, so patch the module itself
        import olefile
        with mock.patch.object(olefile, "isOleFile", lambda p: True), \
                mock.patch.object(olefile, "OleFileIO", lambda p: Ole()):
            meta = ps._extract_ole("anything.doc")
        self.assertEqual(meta["Author"], "str.author")
        self.assertEqual(meta["LastModifiedBy"], "bytes.user")


class TestMetadataSanitizing(unittest.TestCase):
    """Metadata is authored by the target, so it is untrusted input."""

    def test_terminal_escapes_are_stripped(self):
        cleaned = ps._clean_metadata_value("\x1b[2J\x1b[31mred\x1b[0m")
        self.assertNotIn("\x1b", cleaned)

    def test_newlines_cannot_forge_summary_lines(self):
        cleaned = ps._clean_metadata_value("real\n  - fake.admin\n  - fake.root")
        self.assertNotIn("\n", cleaned)

    def test_length_is_capped(self):
        cleaned = ps._clean_metadata_value("x" * 5000)
        self.assertLessEqual(len(cleaned), ps.MAX_METADATA_LENGTH + 3)

    def test_ordinary_values_are_untouched(self):
        self.assertEqual(ps._clean_metadata_value("Jane Doe"), "Jane Doe")
        self.assertEqual(ps._clean_metadata_value("  padded  "), "padded")


# ---------------------------------------------------------------------------
# Output
# ---------------------------------------------------------------------------

class TestOutput(unittest.TestCase):
    def setUp(self):
        self.dir = tempfile.mkdtemp(prefix="pyspyder-out-")
        self.addCleanup(shutil.rmtree, self.dir, ignore_errors=True)
        self.csv = os.path.join(self.dir, "out.csv")

    def test_csv_header_is_the_union_of_all_keys(self):
        results = [
            {"file": "a.pdf", "url": "https://example.com/a.pdf",
             "metadata": {"Author": "one"}},
            {"file": "b.docx", "url": "https://example.com/b.docx",
             "metadata": {"Company": "two", "Manager": "three"}},
        ]
        ps.export_csv(results, self.csv)
        import csv as csv_module
        with open(self.csv, newline="", encoding="utf-8") as fh:
            rows = list(csv_module.DictReader(fh))
        self.assertEqual(set(rows[0].keys()),
                         {"Filename", "URL", "Author", "Company", "Manager"})
        self.assertEqual(rows[0]["Author"], "one")
        self.assertEqual(rows[1]["Manager"], "three")

    def test_formula_injection_is_neutralized(self):
        # An Author field of "=cmd|'/c calc'!A0" is a live payload in Excel.
        hostile = {"Author": "=cmd|'/c calc.exe'!A0", "Title": "+1",
                   "Company": "-2+3", "Keywords": "@SUM(1)"}
        ps.export_csv([{"file": "-evil.pdf", "url": "https://example.com/e.pdf",
                        "metadata": hostile}], self.csv)
        with open(self.csv, encoding="utf-8") as fh:
            body = fh.read().splitlines()[1]
        for dangerous in ["=cmd", "+1", "-2+3", "@SUM"]:
            self.assertNotIn(f",{dangerous}", body, dangerous)
        self.assertIn("'=cmd", body)
        self.assertIn("'-evil.pdf", body)

    def test_summary_aggregates_people_software_and_emails(self):
        results = [
            {"metadata": {"Author": "jsmith", "Creator": "Word",
                          "Subject": "mail Admin@Example.com"}},
            {"metadata": {"Author": "jsmith", "LastModifiedBy": "admin",
                          "Manager": "boss", "Company": "Acme"}},
            {"metadata": {}},
        ]
        summary = ps.format_metadata_summary(results)
        self.assertIn("2 file(s) with metadata, 1 without", summary)
        for expected in ["jsmith", "admin", "boss", "Acme", "Word",
                         "admin@example.com"]:
            self.assertIn(expected, summary)
        self.assertEqual(summary.count("jsmith"), 1, "should be deduplicated")

    def test_summary_handles_no_results(self):
        self.assertEqual(ps.format_metadata_summary([]), "No metadata found.")

    def test_csv_with_no_metadata_at_all(self):
        ps.export_csv([{"file": "a.pdf", "url": "https://example.com/a.pdf",
                        "metadata": {}}], self.csv)
        with open(self.csv, encoding="utf-8") as fh:
            self.assertEqual(fh.read().splitlines()[0], "Filename,URL")

    def test_multiline_metadata_stays_on_one_csv_row(self):
        results = [{"file": "a.pdf", "url": "https://example.com/a.pdf",
                    "metadata": {"Author": ps._clean_metadata_value("a\nb\nc")}}]
        ps.export_csv(results, self.csv)
        with open(self.csv, encoding="utf-8") as fh:
            self.assertEqual(len(fh.read().strip().splitlines()), 2)


class TestDownloadEdges(unittest.TestCase):
    def test_empty_url_list_is_harmless(self):
        out = tempfile.mkdtemp(prefix="pyspyder-empty-")
        self.addCleanup(shutil.rmtree, out, ignore_errors=True)
        self.assertEqual(ps.download_files([], out, FakeSession(), delay=0), [])


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

class TestCli(unittest.TestCase):
    def _rejects(self, argv):
        with open(os.devnull, "w") as devnull:
            stderr, sys.stderr = sys.stderr, devnull
            try:
                with self.assertRaises(SystemExit):
                    ps.parse_args(argv)
            finally:
                sys.stderr = stderr

    def test_a_target_is_required(self):
        self._rejects([])

    def test_target_modes_are_mutually_exclusive(self):
        # Passing two used to be accepted, with one silently ignored.
        self._rejects(["-d", "example.com", "--url-list", "u.txt"])
        self._rejects(["-d", "example.com", "--metadata-only", "some/dir"])
        self._rejects(["--url-list", "u.txt", "--metadata-only", "some/dir"])

    def test_csv_with_download_only_is_rejected(self):
        self._rejects(["-d", "example.com", "--download-only", "--csv", "out.csv"])

    def test_numeric_arguments_are_validated(self):
        self._rejects(["-d", "example.com", "--delay", "-1"])
        self._rejects(["-d", "example.com", "--depth", "-1"])
        self._rejects(["-d", "example.com", "--max-pages", "0"])
        self._rejects(["-d", "example.com", "--max-file-size", "-5"])

    def test_defaults(self):
        args = ps.parse_args(["-d", "example.com"])
        self.assertEqual(args.depth, 2)
        self.assertEqual(args.max_pages, 500)
        self.assertEqual(args.delay, 1.0)
        self.assertEqual(ps.normalize_file_types(args.file_types),
                         ps.DEFAULT_FILE_TYPES)

    def test_file_types_normalization(self):
        self.assertEqual(ps.normalize_file_types("pdf,docx"), ["pdf", "docx"])
        self.assertEqual(ps.normalize_file_types(" .PDF , .DocX "), ["pdf", "docx"])
        self.assertEqual(ps.normalize_file_types("pdf,,."), ["pdf"])

    def test_leading_dot_file_types_still_match_urls(self):
        # -f .pdf is the natural way to type it and used to match nothing.
        self.assertTrue(ps.is_target_file("https://example.com/a.pdf", [".pdf"]))
        self.assertTrue(ps.is_target_file("https://example.com/a.PDF", ["pdf"]))
        self.assertFalse(ps.is_target_file("https://example.com/a.txt", ["pdf"]))


if __name__ == "__main__":
    unittest.main(verbosity=2)
